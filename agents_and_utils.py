# agents_and_utils.py
import os, io
import fitz, docx, csv
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import ollama
from mcp import MCPMessage

model = SentenceTransformer('paraphrase-MiniLM-L3-v2')  # smaller & faster

class FAISSVectorStore:
    def __init__(self):
        self.index = faiss.IndexFlatL2(384)
        self.texts = []

    def add(self, embeddings, texts):
        embeddings = np.array(embeddings)
        if embeddings.ndim == 1:
            embeddings = embeddings.reshape(1, -1)
        embeddings = embeddings.astype("float32")
        self.index.add(embeddings)
        self.texts.extend(texts)

    def search(self, query_emb, k=5):
        if self.index.ntotal == 0:
            return []
        D, I = self.index.search(query_emb, k)
        return [self.texts[i] for i in I[0] if 0 <= i < len(self.texts)]

vector_store = FAISSVectorStore()

def parse_document(file):
    chunks = []
    name = file.name.lower()
    content = file.read()

    if name.endswith(".pdf"):
        with fitz.open(stream=content, filetype="pdf") as doc:
            for page in doc:
                chunks.append(page.get_text())
    elif name.endswith(".docx"):
        doc = docx.Document(io.BytesIO(content))
        chunks = [p.text for p in doc.paragraphs if p.text.strip()]
    elif name.endswith(".txt") or name.endswith(".md"):
        chunks = content.decode(errors="ignore").split("\n\n")
    elif name.endswith(".csv"):
        text = content.decode(errors="ignore").splitlines()
        reader = csv.reader(text)
        chunks = [" | ".join(row) for row in reader]
    elif name.endswith(".pptx"):
        prs = Presentation(io.BytesIO(content))
        for slide in prs.slides:
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            chunks.append(text)
    return [c.strip() for c in chunks if len(c.strip()) > 20][:150]  # limit to 150 chunks

def embed_chunks(chunks):
    if not chunks:
        return np.empty((0, 384), dtype="float32")
    embeddings = model.encode(chunks)
    return np.array(embeddings).astype("float32")

class IngestionAgent:
    def __init__(self, bus): self.bus = bus; bus.register("IngestionAgent", self.handle)
    def handle(self, msg):
        all_chunks = []
        for file in msg.payload["files"]:
            all_chunks.extend(parse_document(file))

        if not all_chunks:
            print("⚠️ No valid content parsed from documents.")
            return

        embeddings = embed_chunks(all_chunks)
        vector_store.add(embeddings, all_chunks)
        new_msg = MCPMessage("IngestionAgent", "RetrievalAgent", "INGESTION_DONE", msg.trace_id, {"query": msg.payload["query"]})
        self.bus.send(new_msg)

class RetrievalAgent:
    def __init__(self, bus): self.bus = bus; bus.register("RetrievalAgent", self.handle)
    def handle(self, msg):
        q_emb = embed_chunks([msg.payload["query"]])
        if q_emb.shape[0] == 0:
            return
        top_chunks = vector_store.search(q_emb, k=5)
        if not top_chunks:
            print("⚠️ No chunks retrieved. Possible empty index.")
            return
        new_msg = MCPMessage("RetrievalAgent", "LLMResponseAgent", "RETRIEVAL_RESULT", msg.trace_id, {
            "top_chunks": top_chunks,
            "query": msg.payload["query"]
        })
        self.bus.send(new_msg)

class LLMResponseAgent:
    def __init__(self, bus, model_name):
        self.bus = bus
        self.model = model_name
        self.final_answer = ""
        self.final_context = []
        self.bus.register("LLMResponseAgent", self.handle)

    def handle(self, msg):
        context_chunks = msg.payload.get("top_chunks", [])
        user_query = msg.payload.get("query", "")

        if not context_chunks:
            self.final_answer = "⚠️ I couldn't extract meaningful content from the uploaded document."
            self.final_context = []
            return

        system_prompt = "You are a helpful assistant. Use the provided context to answer the question."
        context = "\n---\n".join(context_chunks)
        prompt = f"{system_prompt}\nContext:\n{context}\n\nQuestion: {user_query}\nAnswer:"

        response = ollama.chat(
            model=self.model,
            messages=[{"role": "user", "content": prompt}]
        )
        self.final_answer = response['message']['content']
        self.final_context = context_chunks